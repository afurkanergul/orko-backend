# backend/app/integrations/files/parsers/fallback.py
"""
Fallback parser — replaces textract safely
-------------------------------------------
Handles multiple document types without external dependencies.
Compatible with Render and Python 3.11+.
"""

import os
from PyPDF2 import PdfReader
from docx import Document


def read_with_textract(path: str) -> str:
    """
    Fallback reader (safe replacement for textract).
    Supports PDF, DOCX, TXT, CSV, LOG, and PPTX files.
    Returns text content or empty string if unreadable.
    """
    ext = os.path.splitext(path)[1].lower()

    try:
        # --- PDF ---
        if ext == ".pdf":
            reader = PdfReader(path)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            return text.strip()

        # --- DOCX ---
        elif ext == ".docx":
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs).strip()

        # --- Plain text / CSV / logs ---
        elif ext in (".txt", ".csv", ".log"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read().strip()

        # --- PowerPoint ---
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            slides = [
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            ]
            return "\n".join(slides).strip()

        # --- Unsupported file types ---
        else:
            print(f"[WARN] Unsupported file type: {ext}")
            return "[Unsupported file type or binary data]"

    except Exception as e:
        print(f"[ERROR] Fallback extraction failed for {path}: {e}")
        return ""
